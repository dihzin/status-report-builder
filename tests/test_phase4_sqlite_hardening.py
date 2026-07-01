import asyncio
import base64
import gc
import importlib
import sqlite3
import time
from pathlib import Path

from pptx import Presentation

from backend.db import ensure_db
from backend.excel_reader import create_template
from backend.exporter import export_pdf, export_pptx
from backend.services.report_service import ReportService


def _root(tmp_path: Path) -> Path:
    root = tmp_path / "proj"
    root.mkdir(parents=True, exist_ok=True)
    return root


def _db_path(root: Path) -> Path:
    return root / "data" / "status_builder.db"


def _excel_path(root: Path) -> Path:
    return root / "status_projeto.xlsx"


def _service(root: Path) -> ReportService:
    return ReportService(root)


def _patch_db(monkeypatch, root: Path):
    import backend.db as db_module

    monkeypatch.setattr(db_module, "ROOT_DIR", root)
    monkeypatch.setattr(db_module, "DATA_DIR", root / "data")
    monkeypatch.setattr(db_module, "DB_PATH", root / "data" / "status_builder.db")


def _current_snapshot_meta(db_path: Path):
    conn = sqlite3.connect(str(db_path))
    cur = conn.cursor()
    cur.execute("SELECT source, version_number FROM report_snapshots WHERE project_key='default' AND is_current=1")
    row = cur.fetchone()
    conn.close()
    return row


def test_first_run_without_db_with_excel_seeds_sqlite(tmp_path, monkeypatch):
    root = _root(tmp_path)
    _patch_db(monkeypatch, root)
    create_template(str(_excel_path(root)))

    svc = _service(root)
    payload = svc.get_status_payload()

    assert payload["reportData"]
    assert _db_path(root).exists()
    src, version = _current_snapshot_meta(_db_path(root))
    assert src == "excel_import"
    assert version == 1


def test_existing_db_without_excel_still_runs(tmp_path, monkeypatch):
    root = _root(tmp_path)
    _patch_db(monkeypatch, root)
    excel = _excel_path(root)
    create_template(str(excel))

    svc = _service(root)
    svc.get_status_payload()
    gc.collect()
    for _ in range(5):
        try:
            excel.unlink()
            break
        except PermissionError:
            time.sleep(0.1)
    else:
        raise AssertionError("Nao foi possivel remover status_projeto.xlsx durante o teste")

    payload = svc.get_status_payload()
    assert payload["reportData"]


def test_empty_db_without_excel_creates_template_and_seeds_sqlite(tmp_path, monkeypatch):
    root = _root(tmp_path)
    _patch_db(monkeypatch, root)
    ensure_db()

    svc = _service(root)
    payload = svc.get_status_payload()

    assert payload["reportData"]
    assert _excel_path(root).exists()
    src, version = _current_snapshot_meta(_db_path(root))
    assert src == "excel_import"
    assert version == 1


def test_save_legacy_and_canonical_with_unique_versions_and_single_current(tmp_path, monkeypatch):
    root = _root(tmp_path)
    _patch_db(monkeypatch, root)
    create_template(str(_excel_path(root)))
    svc = _service(root)
    svc.get_status_payload()

    svc.save_payload({"config": {"project_subtitle": "LEGACY_SUB"}})
    svc.save_payload({"reportData": {"config": {"project_name": "CANON_NAME"}}})

    data = svc.get_status_payload()
    assert data["reportData"]["config"]["project_name"] == "CANON_NAME"
    assert data["reportData"]["config"]["project_subtitle"] == "LEGACY_SUB"

    conn = sqlite3.connect(str(_db_path(root)))
    cur = conn.cursor()
    cur.execute("SELECT COUNT(*) FROM report_snapshots WHERE project_key='default' AND is_current=1")
    current_count = cur.fetchone()[0]
    cur.execute("SELECT COUNT(*) FROM (SELECT version_number, COUNT(*) c FROM report_snapshots WHERE project_key='default' GROUP BY version_number HAVING c > 1)")
    duplicated_versions = cur.fetchone()[0]
    conn.close()

    assert current_count == 1
    assert duplicated_versions == 0


def test_excel_manual_change_does_not_override_sqlite(tmp_path, monkeypatch):
    root = _root(tmp_path)
    _patch_db(monkeypatch, root)
    excel = _excel_path(root)
    create_template(str(excel))
    svc = _service(root)
    svc.get_status_payload()
    svc.save_payload({"reportData": {"config": {"project_name": "SQLITE_PRIORITY"}}})

    import openpyxl

    wb = openpyxl.load_workbook(str(excel))
    ws = wb["CONFIG"]
    for row in ws.iter_rows(min_row=1):
        if row[0].value == "project_name":
            row[1].value = "EXCEL_OVERRIDE_ATTEMPT"
            break
    wb.save(str(excel))
    wb.close()

    payload = svc.get_status_payload()
    assert payload["reportData"]["config"]["project_name"] == "SQLITE_PRIORITY"


def test_export_pdf_and_pptx_use_sqlite_data(tmp_path, monkeypatch):
    root = _root(tmp_path)
    _patch_db(monkeypatch, root)
    create_template(str(_excel_path(root)))
    svc = _service(root)
    svc.get_status_payload()
    svc.save_payload({"reportData": {"config": {"project_name": "PPTX_SQLITE_NAME"}}})

    import backend.exporter as exporter_module

    monkeypatch.setattr(exporter_module, "ROOT_DIR", root)
    monkeypatch.setattr(exporter_module, "PDF_DIR", root / "exports" / "pdf")
    monkeypatch.setattr(exporter_module, "PPTX_DIR", root / "exports" / "pptx")
    exporter_module.PDF_DIR.mkdir(parents=True, exist_ok=True)
    exporter_module.PPTX_DIR.mkdir(parents=True, exist_ok=True)

    def fake_worker(cmd: dict, timeout: int = 60):
        p = Path(cmd["path"])
        p.parent.mkdir(parents=True, exist_ok=True)
        action = cmd.get("action")
        if action == "pdf":
            p.write_bytes(b"%PDF-FAKE")
        elif action == "screenshot":
            png_b64 = "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO7Zl9sAAAAASUVORK5CYII="
            p.write_bytes(base64.b64decode(png_b64))

    monkeypatch.setattr(exporter_module, "_run_worker", fake_worker)

    pdf_path = asyncio.run(export_pdf("http://127.0.0.1:8000"))
    pptx_path = asyncio.run(export_pptx("http://127.0.0.1:8000", data_provider=svc.get_current_report_data))

    assert Path(pdf_path).exists()
    assert Path(pptx_path).exists()

    prs = Presentation(pptx_path)
    slide1_text = []
    for shape in prs.slides[0].shapes:
        if hasattr(shape, "text"):
            slide1_text.append(shape.text)
    assert any("PPTX_SQLITE_NAME" in t for t in slide1_text)


def test_hardening_indexes_are_created(tmp_path, monkeypatch):
    root = _root(tmp_path)
    _patch_db(monkeypatch, root)
    ensure_db()

    conn = sqlite3.connect(str(root / "data" / "status_builder.db"))
    cur = conn.cursor()
    cur.execute("PRAGMA index_list('report_snapshots')")
    indexes = {r[1] for r in cur.fetchall()}
    conn.close()

    assert "uq_report_snapshots_project_version" in indexes
    assert "uq_report_snapshots_single_current" in indexes


def test_legacy_duplicate_versions_are_sanitized_before_unique_index(tmp_path, monkeypatch):
    root = _root(tmp_path)
    _patch_db(monkeypatch, root)
    db_path = root / "data" / "status_builder.db"
    db_path.parent.mkdir(parents=True, exist_ok=True)

    conn = sqlite3.connect(str(db_path))
    cur = conn.cursor()
    cur.execute(
        """
        CREATE TABLE report_snapshots (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            project_key TEXT NOT NULL,
            report_name TEXT,
            report_date TEXT,
            version_number INTEGER NOT NULL DEFAULT 1,
            is_current INTEGER NOT NULL DEFAULT 1,
            source TEXT NOT NULL,
            report_data_json TEXT NOT NULL,
            legacy_data_json TEXT NOT NULL,
            created_at TEXT NOT NULL DEFAULT CURRENT_TIMESTAMP,
            updated_at TEXT NOT NULL DEFAULT CURRENT_TIMESTAMP
        )
        """
    )
    rows = [
        ("legacy", 1, 0, "legacy_seed", "2026-01-01 10:00:00"),
        ("legacy", 1, 1, "legacy_seed", "2026-01-02 10:00:00"),
        ("legacy", 1, 0, "legacy_seed", "2026-01-03 10:00:00"),
        ("legacy", 2, 0, "legacy_seed", "2026-01-04 10:00:00"),
    ]
    for project_key, version_number, is_current, source, created_at in rows:
        cur.execute(
            """
            INSERT INTO report_snapshots(
                project_key, report_name, report_date, version_number, is_current, source,
                report_data_json, legacy_data_json, created_at, updated_at
            ) VALUES (?, 'R', '2026-01-01', ?, ?, ?, '{}', '{}', ?, ?)
            """,
            (project_key, version_number, is_current, source, created_at, created_at),
        )
    conn.commit()
    conn.close()

    ensure_db()

    conn = sqlite3.connect(str(db_path))
    cur = conn.cursor()
    cur.execute("SELECT COUNT(*) FROM report_snapshots WHERE project_key='legacy'")
    total_rows = cur.fetchone()[0]
    cur.execute(
        """
        SELECT COUNT(*) FROM (
            SELECT version_number, COUNT(*) c
            FROM report_snapshots
            WHERE project_key='legacy'
            GROUP BY version_number
            HAVING c > 1
        )
        """
    )
    duplicated_versions = cur.fetchone()[0]
    cur.execute("SELECT COUNT(*) FROM report_snapshots WHERE project_key='legacy' AND is_current=1")
    current_count = cur.fetchone()[0]
    cur.execute("SELECT COUNT(*) FROM report_snapshots WHERE project_key='legacy' AND version_number=1 AND is_current=1")
    preserved_current = cur.fetchone()[0]
    conn.close()

    assert total_rows == 4
    assert duplicated_versions == 0
    assert current_count == 1
    assert preserved_current == 1


def test_legacy_dedup_migration_is_idempotent(tmp_path, monkeypatch):
    root = _root(tmp_path)
    _patch_db(monkeypatch, root)
    db_path = root / "data" / "status_builder.db"
    db_path.parent.mkdir(parents=True, exist_ok=True)

    conn = sqlite3.connect(str(db_path))
    cur = conn.cursor()
    cur.execute(
        """
        CREATE TABLE report_snapshots (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            project_key TEXT NOT NULL,
            report_name TEXT,
            report_date TEXT,
            version_number INTEGER NOT NULL DEFAULT 1,
            is_current INTEGER NOT NULL DEFAULT 1,
            source TEXT NOT NULL,
            report_data_json TEXT NOT NULL,
            legacy_data_json TEXT NOT NULL,
            created_at TEXT NOT NULL DEFAULT CURRENT_TIMESTAMP,
            updated_at TEXT NOT NULL DEFAULT CURRENT_TIMESTAMP
        )
        """
    )
    cur.executemany(
        """
        INSERT INTO report_snapshots(
            project_key, report_name, report_date, version_number, is_current, source,
            report_data_json, legacy_data_json, created_at, updated_at
        ) VALUES ('legacy', 'R', '2026-01-01', ?, ?, 'legacy_seed', '{}', '{}', ?, ?)
        """,
        [
            (3, 0, "2026-01-01 10:00:00", "2026-01-01 10:00:00"),
            (3, 1, "2026-01-02 10:00:00", "2026-01-02 10:00:00"),
            (3, 0, "2026-01-03 10:00:00", "2026-01-03 10:00:00"),
        ],
    )
    conn.commit()
    conn.close()

    ensure_db()
    conn = sqlite3.connect(str(db_path))
    cur = conn.cursor()
    cur.execute("SELECT id, version_number, is_current FROM report_snapshots WHERE project_key='legacy' ORDER BY id")
    first_run_rows = cur.fetchall()
    conn.close()

    ensure_db()
    conn = sqlite3.connect(str(db_path))
    cur = conn.cursor()
    cur.execute("SELECT id, version_number, is_current FROM report_snapshots WHERE project_key='legacy' ORDER BY id")
    second_run_rows = cur.fetchall()
    cur.execute(
        """
        SELECT COUNT(*) FROM (
            SELECT version_number, COUNT(*) c
            FROM report_snapshots
            WHERE project_key='legacy'
            GROUP BY version_number
            HAVING c > 1
        )
        """
    )
    duplicated_versions = cur.fetchone()[0]
    conn.close()

    assert first_run_rows == second_run_rows
    assert duplicated_versions == 0


def test_status_does_not_depend_on_excel_validation_when_flag_false(tmp_path, monkeypatch):
    root = _root(tmp_path)
    _patch_db(monkeypatch, root)
    create_template(str(_excel_path(root)))

    import backend.services.report_service as report_service_module

    monkeypatch.setenv("VALIDATE_EXCEL_SCHEMA", "false")
    importlib.reload(report_service_module)

    calls = {"count": 0}

    def fake_validate_schema(_):
        calls["count"] += 1
        raise AssertionError("validate_schema nao deveria ser chamado com VALIDATE_EXCEL_SCHEMA=false")

    monkeypatch.setattr(report_service_module, "validate_schema", fake_validate_schema)
    svc = report_service_module.ReportService(root)
    payload = svc.get_status_payload()

    assert payload["reportData"]
    assert payload["validation_errors"] == []
    assert calls["count"] == 0
