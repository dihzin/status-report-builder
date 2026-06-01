import asyncio
import json
import subprocess
import sys
import tempfile
from concurrent.futures import ThreadPoolExecutor
from datetime import datetime
from pathlib import Path
from typing import Callable

_WORKER = Path(__file__).parent / "export_worker.py"
_executor = ThreadPoolExecutor(max_workers=2)


def _run_worker(cmd: dict, timeout: int = 60) -> None:
    result = subprocess.run(
        [sys.executable, str(_WORKER), json.dumps(cmd)],
        capture_output=True,
        text=True,
        timeout=timeout,
    )
    if result.returncode != 0:
        raise RuntimeError(result.stderr or result.stdout or "export_worker falhou")

ROOT_DIR = Path(__file__).resolve().parent.parent
PDF_DIR  = ROOT_DIR / "exports" / "pdf"
PPTX_DIR = ROOT_DIR / "exports" / "pptx"
PDF_DIR.mkdir(parents=True, exist_ok=True)
PPTX_DIR.mkdir(parents=True, exist_ok=True)


def _stamped(stem: str, ext: str) -> str:
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    return f"{stem}_{ts}{ext}"


def _clean(v, default=""):
    if v is None:
        return default
    s = str(v).strip()
    return s if s else default


def _token(pres, key, default=""):
    v = pres.get(key)
    return _clean(v, default) if v is not None else default


def _hex_to_rgb_tuple(hex_color: str, fallback=(42, 114, 73)):
    s = _clean(hex_color, "").lstrip("#")
    if len(s) != 6:
        return fallback
    try:
        return tuple(int(s[i:i + 2], 16) for i in (0, 2, 4))
    except Exception:
        return fallback


async def export_pdf(frontend_url: str) -> str:
    pdf_path = PDF_DIR / _stamped("status_report", ".pdf")
    cmd = {"action": "pdf", "url": frontend_url, "path": str(pdf_path)}
    loop = asyncio.get_event_loop()
    await loop.run_in_executor(_executor, _run_worker, cmd)
    return str(pdf_path)


async def export_pptx(frontend_url: str, data_provider: Callable[[], dict] | None = None) -> str:
    from pptx import Presentation
    from pptx.util import Inches
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

    pptx_path = PPTX_DIR / _stamped("status_report", ".pptx")
    data = data_provider() if callable(data_provider) else {}
    data = data or {}
    cfg = data.get("config", {}) or {}
    branding = data.get("branding", {}) or {}
    pres = data.get("presentation_config", {}) or {}
    rodape = data.get("rodape", {}) or {}
    fases = data.get("fases", []) or []
    marcos = data.get("marcos", []) or []
    pendencias = data.get("pendencias_criticas", []) or []
    acoes = data.get("proximas_acoes", []) or []

    screenshot_path = Path(tempfile.gettempdir()) / "status_screenshot.png"
    cmd = {"action": "screenshot", "url": frontend_url, "path": str(screenshot_path)}
    loop = asyncio.get_event_loop()
    await loop.run_in_executor(_executor, _run_worker, cmd)

    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)

    primary_rgb = _hex_to_rgb_tuple(branding.get("cor_primaria", "#2a7249"), fallback=(42, 114, 73))
    secondary_rgb = _hex_to_rgb_tuple(branding.get("cor_secundaria", "#1a4f35"), fallback=(26, 79, 53))
    white = RGBColor(255, 255, 255)
    dark = RGBColor(*_hex_to_rgb_tuple(_token(pres, "text_on_light_primary", "#1E2228"), fallback=(30, 34, 40)))
    gray = RGBColor(*_hex_to_rgb_tuple(_token(pres, "text_on_light_secondary", "#454B54"), fallback=(106, 112, 122)))
    primary = RGBColor(*primary_rgb)
    secondary = RGBColor(*secondary_rgb)
    slide_simple_bg = RGBColor(*_hex_to_rgb_tuple(_token(pres, "slide_simple_bg", "#F8FBFD"), fallback=(248, 250, 253)))
    font_family = _token(pres, "font_family", "Inter")

    def add_bg(slide, rgb):
        bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = rgb
        bg.line.fill.background()
        return bg

    def add_text(slide, text, left, top, width, height, size=20, bold=False, color=dark, align=PP_ALIGN.LEFT):
        tx = slide.shapes.add_textbox(left, top, width, height)
        tf = tx.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = _clean(text, "")
        font = run.font
        font.name = font_family
        font.size = Pt(size)
        font.bold = bold
        font.color.rgb = color
        return tx

    # Slide 1: Capa dinâmica
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s1, secondary)
    band = s1.shapes.add_shape(1, Inches(0), Inches(5.95), prs.slide_width, Inches(1.55))
    band.fill.solid()
    band.fill.fore_color.rgb = primary
    band.line.fill.background()

    add_text(
        s1,
        _clean(cfg.get("report_title"), "STATUS REPORT"),
        Inches(0.85), Inches(1.0), Inches(8.8), Inches(0.75),
        size=34, bold=True, color=white
    )
    add_text(
        s1,
        _clean(cfg.get("project_name"), "Projeto"),
        Inches(0.85), Inches(1.9), Inches(9.2), Inches(0.6),
        size=21, bold=True, color=white
    )
    add_text(
        s1,
        _clean(cfg.get("project_subtitle"), ""),
        Inches(0.85), Inches(2.45), Inches(10.2), Inches(0.7),
        size=14, color=RGBColor(232, 240, 255)
    )
    add_text(
        s1,
        f"Cliente: {_clean(cfg.get('sponsor'), '-')}",
        Inches(0.85), Inches(6.25), Inches(4.4), Inches(0.35),
        size=12, bold=True, color=white
    )
    add_text(
        s1,
        f"Apresentador: {_clean(cfg.get('owner_name'), '-')}",
        Inches(4.75), Inches(6.25), Inches(4.7), Inches(0.35),
        size=12, bold=True, color=white
    )
    add_text(
        s1,
        f"Data: {_clean(cfg.get('report_date'), '-')}",
        Inches(9.2), Inches(6.25), Inches(3.0), Inches(0.35),
        size=12, bold=True, color=white
    )

    logo_path = _clean(branding.get("logo_path"), "")
    if logo_path:
        abs_logo = (ROOT_DIR / "frontend" / logo_path).resolve() if not Path(logo_path).is_absolute() else Path(logo_path)
        if abs_logo.exists():
            s1.shapes.add_picture(str(abs_logo), Inches(10.35), Inches(0.7), Inches(2.0), Inches(1.45))

    # Slide 2: Onepage (screenshot)
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    s2.shapes.add_picture(str(screenshot_path), 0, 0, prs.slide_width, prs.slide_height)

    # Slide 3: Cronograma e Marcos (Gantt)
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s3, slide_simple_bg)
    add_text(s3, "Cronograma Executivo", Inches(0.7), Inches(0.35), Inches(8.0), Inches(0.6), size=28, bold=True, color=secondary)
    add_text(
        s3,
        f"Fase atual: {_clean(cfg.get('current_phase'), '-')}  |  Dia {_clean(cfg.get('current_day'), '-')} de {_clean(cfg.get('total_days'), '-')}",
        Inches(0.72), Inches(0.95), Inches(10.8), Inches(0.35), size=12, color=gray
    )

    # Preferir dados GANTT; fallback para FASES/MARCOS legado
    gantt_tasks = data.get("gantt_tarefas") or []
    gantt_marcos_list = data.get("gantt_marcos") or []
    has_gantt = bool(gantt_tasks or gantt_marcos_list)

    top = 1.45
    row_h = 0.42
    if has_gantt:
        add_text(s3, "Tarefas", Inches(0.72), Inches(top), Inches(5.8), Inches(0.35), size=13, bold=True, color=primary)
        for i, t in enumerate(gantt_tasks[:8], 1):
            line = f"{i}. {_clean(t.get('nome'), '-')}  |  {_clean(t.get('inicio'), '-')} → {_clean(t.get('fim'), '-')}  |  {_clean(t.get('status'), '-')}"
            add_text(s3, line, Inches(0.72), Inches(top + i * row_h), Inches(6.2), Inches(0.35), size=11, color=dark)

        mtop = 1.45
        add_text(s3, "Marcos", Inches(7.0), Inches(mtop), Inches(5.2), Inches(0.35), size=13, bold=True, color=primary)
        for i, m in enumerate(gantt_marcos_list[:8], 1):
            line = f"{i}. {_clean(m.get('nome'), '-')}  |  {_clean(m.get('data'), '-')}  |  {_clean(m.get('status'), '-')}"
            add_text(s3, line, Inches(7.0), Inches(mtop + i * row_h), Inches(5.5), Inches(0.35), size=11, color=dark)
    else:
        add_text(s3, "Fases", Inches(0.72), Inches(top), Inches(5.8), Inches(0.35), size=13, bold=True, color=primary)
        for i, fase in enumerate(fases[:8], 1):
            line = f"{i}. {_clean(fase.get('nome'), '-')}  |  {_clean(fase.get('status'), '-')}  |  {_clean(fase.get('data_alvo'), '-')}"
            add_text(s3, line, Inches(0.72), Inches(top + i * row_h), Inches(6.2), Inches(0.35), size=11, color=dark)

        mtop = 1.45
        add_text(s3, "Marcos", Inches(7.0), Inches(mtop), Inches(5.2), Inches(0.35), size=13, bold=True, color=primary)
        for i, m in enumerate(marcos[:8], 1):
            line = f"{i}. {_clean(m.get('nome'), '-')}  |  {_clean(m.get('status'), '-')}  |  {_clean(m.get('data_alvo'), '-')}"
            add_text(s3, line, Inches(7.0), Inches(mtop + i * row_h), Inches(5.5), Inches(0.35), size=11, color=dark)

    # Slide 4: Detalhamento de riscos e ações
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s4, slide_simple_bg)
    add_text(s4, "Detalhamento: Pendências e Próximas Ações", Inches(0.7), Inches(0.35), Inches(11.8), Inches(0.6), size=24, bold=True, color=secondary)

    add_text(s4, "Pendências Críticas", Inches(0.72), Inches(1.05), Inches(5.8), Inches(0.35), size=13, bold=True, color=primary)
    for i, p in enumerate(pendencias[:7], 1):
        line = f"{_clean(p.get('prioridade'), 'P?')} | {_clean(p.get('item'), '-')}"
        meta = f"{_clean(p.get('responsaveis'), '-')} | {_clean(p.get('status'), '-')}"
        add_text(s4, line, Inches(0.72), Inches(1.05 + i * 0.55), Inches(5.9), Inches(0.28), size=11, bold=True, color=dark)
        add_text(s4, meta, Inches(0.72), Inches(1.22 + i * 0.55), Inches(5.9), Inches(0.24), size=10, color=gray)

    add_text(s4, "Próximas Ações", Inches(7.0), Inches(1.05), Inches(5.2), Inches(0.35), size=13, bold=True, color=primary)
    for i, a in enumerate(acoes[:10], 1):
        line = f"{i}. {_clean(a.get('texto'), '-')}"
        add_text(s4, line, Inches(7.0), Inches(1.05 + i * 0.46), Inches(5.6), Inches(0.34), size=11, color=dark)

    # Slide 5: Encerramento
    s5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s5, secondary)
    add_text(s5, _clean(cfg.get("report_title"), "STATUS REPORT"), Inches(0.8), Inches(1.4), Inches(8.0), Inches(0.8), size=36, bold=True, color=white)
    add_text(s5, "Obrigado.", Inches(0.8), Inches(2.5), Inches(5.0), Inches(0.7), size=32, bold=True, color=white)
    light_text = RGBColor(*_hex_to_rgb_tuple(_token(pres, "text_on_dark_secondary", "rgba(235,242,255)").replace("rgba(", "#").replace(")", "").replace(",", ""), fallback=(235, 242, 255)))
    add_text(
        s5,
        f"Próximo marco: {_clean(rodape.get('milestone_alvo'), _clean(cfg.get('current_phase'), '-'))}",
        Inches(0.8), Inches(3.45), Inches(10.5), Inches(0.55), size=18, color=light_text
    )
    add_text(
        s5,
        f"Data alvo: {_clean(rodape.get('data_alvo'), '-')}  |  Go-Live: {_clean(rodape.get('go_live_previsto'), '-')}",
        Inches(0.8), Inches(4.05), Inches(10.5), Inches(0.45), size=14, color=light_text
    )

    prs.save(str(pptx_path))

    if screenshot_path.exists():
        screenshot_path.unlink()

    return str(pptx_path)
